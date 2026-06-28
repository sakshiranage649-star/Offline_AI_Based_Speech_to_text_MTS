import subprocess
import sys
import os

def install(package):
    subprocess.check_call([sys.executable, "-m", "pip", "install", package])

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
except ImportError:
    print("Installing python-pptx...")
    install("python-pptx")
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()
    
    # -------------------------------------------------------------
    # Slide 1: Title Slide
    # -------------------------------------------------------------
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Multilingual Translation - Project Architecture"
    subtitle.text = "Offline Multilingual AI Translation Web App\nProject Overview & Architecture flow"
    
    # -------------------------------------------------------------
    # Slide 2: Technology Stack Used
    # -------------------------------------------------------------
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "Technology Stack Used"
    
    tf = body_shape.text_frame
    tf.text = "Backend (Server & AI Processing):"
    
    p = tf.add_paragraph()
    p.text = "Python (FastAPI, Uvicorn) for the Web Server."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "HuggingFace Transformers & PyTorch (AI Translation)."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Vosk (Offline Speech-to-Text Engine)."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "PyPDF2 & MoviePy (File audio/text extraction)."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Database:"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "SQLite (Stores history offline persistently)."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Frontend (User Interface):"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "HTML, CSS, JavaScript (Served statically by FastAPI)."
    p.level = 1
    
    # -------------------------------------------------------------
    # Slide 3: How the Application Works
    # -------------------------------------------------------------
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "How the Website Works (Step-by-step)"
    
    tf = body_shape.text_frame
    tf.text = "1. User Input: User types text or uploads a file (PDF/Video/Audio)."
    
    p = tf.add_paragraph()
    p.text = "2. Data Extraction:"
    
    sp1 = tf.add_paragraph()
    sp1.text = "PDF text is read via PyPDF2."
    sp1.level = 1
    
    sp2 = tf.add_paragraph()
    sp2.text = "A/V files are passed to Vosk for STT transcription."
    sp2.level = 1
    
    p = tf.add_paragraph()
    p.text = "3. AI Translation: Text is sent to Transformers locally for translation."
    
    p = tf.add_paragraph()
    p.text = "4. Database Storage: Input, Output, and Languages are stored in SQLite."
    
    p = tf.add_paragraph()
    p.text = "5. Display Result: FastAPI returns translated text to simple HTML/JS frontend."
    
    # -------------------------------------------------------------
    # Slide 4: Graphical Workflow / Architecture Box (Blank for user to modify)
    # -------------------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[5]) # Title only layout
    shapes = slide.shapes
    title_shape = shapes.title
    title_shape.text = "Architecture Workflow Flowchart"
    
    # Let's add some basic blocks so the user can just restyle them
    left_margin = Inches(0.5)
    
    # User Box
    shapes.add_shape(1, left_margin, Inches(2), Inches(1.5), Inches(1)).text_frame.text = "User Input (Upload / Type)"
    
    # Frontend Box
    shapes.add_shape(1, Inches(2.5), Inches(2), Inches(1.5), Inches(1)).text_frame.text = "Frontend UI (HTML/JS)"
    
    # Backend Box
    shapes.add_shape(1, Inches(4.5), Inches(2), Inches(1.8), Inches(1)).text_frame.text = "Backend Server (FastAPI)"
    
    # Utilities Box
    shapes.add_shape(1, Inches(4.5), Inches(3.5), Inches(1.8), Inches(1)).text_frame.text = "Data Extractors (PDF/Vosk)"
    
    # Translators Box
    shapes.add_shape(1, Inches(7.0), Inches(2), Inches(2), Inches(1)).text_frame.text = "AI Translation Engine (Transformers)"
    
    # DB Box
    shapes.add_shape(1, Inches(4.5), Inches(5.0), Inches(1.8), Inches(1)).text_frame.text = "SQLite Database"
    
    save_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "Multilingual_Translation_Architecture_Presentation.pptx")
    prs.save(save_path)
    print(f"Presentation generated and saved to: {save_path}")

if __name__ == "__main__":
    create_presentation()
